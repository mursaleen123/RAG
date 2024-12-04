# Standard Libraries
import datetime

# Django
from django.conf import settings


# Office libraries
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Utils
#from MVPautomation.sharepoint.utils.delete_functions import get_slide_index


def create_gantt(presentation, tasks):
    '''
    Function to create a Gantt chart in a PowerPoint slide based on a list of tasks.
    '''

    try:
        # Add a blank slide layout
        slide_index = get_slide_index(presentation, ["[Timeline]"])

        # Access the existing slide or add a new one if necessary
        slide = presentation.slides[slide_index[0]]

        # Optionally, remove all shapes from the new slide
        shapes_to_remove = [shape for shape in slide.shapes]

        # Iterate through and remove each shape
        for shape in shapes_to_remove:
            sp = shape._element
            sp.getparent().remove(sp)


        start_date = min(task.start_date for task in tasks)
        end_date = max(task.end_date for task in tasks)
        total_days = (end_date - start_date).days + 1

        # Define slide width and margins
        slide_width = presentation.slide_width.inches
        left_margin = 2.5
        right_margin = 1.0
        gantt_width = slide_width - left_margin - right_margin

        # Calculate day width dynamically to fit the entire Gantt chart in the slide
        day_width = gantt_width / total_days


        # Add week labels dynamically based on the number of days per week
        current_date = start_date
        top_margin = 2
        label_top_margin = 1
        week_counter = 1
        label_height = Inches(0.4)

        #create the index [Timeline] in a new text frame in a corner top right
        textbox = slide.shapes.add_textbox(Inches(12.2), Inches(0.2), Inches(1.5), Inches(0.5))  # x = 8.5 to move it to the right
        text_frame = textbox.text_frame
        text_frame.text = "[Timeline]"
        text_frame.paragraphs[0].font.size = Pt(10)

        # Add a title to the slide
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(4), Inches(0.5))
        title_text_frame = title_box.text_frame
        title_text_frame.text = "Project Timeline (includes weekly check-ins)"
        title_text_frame.paragraphs[0].font.size = Pt(18)
        title_text_frame.paragraphs[0].font.color.rgb = RGBColor(98, 68, 210)  # Purple color


        # Add week labels with purple boxes behind them
        while current_date <= end_date:
            # Calculate the start and end of the current week (Monday to Sunday)
            week_start = current_date
            week_end = min(week_start + datetime.timedelta(days=6), end_date)  # Ensure the end doesn't go beyond the total range

            # Calculate the width of the week label based on the number of days in this week
            week_days = days_between(week_start, week_end) + 1
            week_label_width = Inches(week_days * day_width)

            # Add a purple box behind the week label
            left = Inches(left_margin + days_between(start_date, week_start) * day_width)
            purple_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(label_top_margin), week_label_width, label_height)
            purple_box.fill.solid()
            purple_box.fill.fore_color.rgb = RGBColor(79, 55, 170)

            # Add the week label
            week_label_box = slide.shapes.add_textbox(left, Inches(label_top_margin), week_label_width, label_height)
            week_label_text_frame = week_label_box.text_frame
            week_label_text_frame.text = f"W{week_counter}"

            # Format the week label
            for paragraph in week_label_text_frame.paragraphs:
                paragraph.font.bold = True
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = RGBColor(255, 255, 255)

            # Move to the next week
            current_date = week_end + datetime.timedelta(days=1)
            week_counter += 1

        # Add tasks as rectangles and add task names as a column on the left
        for i, task in enumerate(tasks):
            # Add task name in the left column
            top = Inches(top_margin + i * 0.7)

        # Calculate task position and width for the Gantt chart
            left = Inches(left_margin + days_between(start_date, task.start_date) * day_width )
            width = Inches(days_between(task.start_date, task.end_date) * day_width)
            height = Inches(0.3)

            # Add task rectangle to the Gantt chart
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(68, 46, 142)
            shape.line.color.rgb = RGBColor(0, 0, 0)
            shape.adjustments[0] = 0.5

            clean_text = task.activity.name.replace("\u200b", "").replace("\n", " ").strip()
            # if the text is too long, put more on top
            if len(clean_text) > 15:
                top -= Inches(0.1)


            task_name_box = slide.shapes.add_textbox(Inches(1), top, Inches(1.5), Inches(0.5))  # Adjust height for better visibility

            # Configure the text box
            text_frame = task_name_box.text_frame
            text_frame.word_wrap = True
            text_frame.text = clean_text
            text_frame.paragraphs[0].font.size = Pt(12)
            text_frame.paragraphs[0].font.bold = True

        print("Gantt chart created successfully!")

    except Exception as e:
        print(f"Error in Gantt population: {e}")
        # raise Exception(f"Error in Gantt population: {e}")


def days_between(d1, d2):
    '''
    Function to calculate the number of days between two dates.
    '''
    return (d2 - d1).days
