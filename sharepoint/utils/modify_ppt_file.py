# Standard Libraries
import itertools

# Django
from django.conf import settings


# Models
#from MVPautomation.project.models import MethodologyModel, ProposalsModel, UseCasesModel

# Office Libraries
from pptx.util import Pt
from pptx.dml.color import RGBColor

# Utils
#from MVPautomation.sharepoint.utils.include_data_json import search_and_include_info


def modify_sp_file(presentation, new_content, methodology):
    '''
    Function to modify a file from Sharepoint in Memory, replace information and return updated file.

    presentation: PPTX presentation File in memory buffer
    new_content: The dict of content to replace in the file
    methodology: Id instance for Methodology
    '''

    try:
        method_data = MethodologyModel.objects.get(id=methodology)

        goals = method_data.specific_goals

        num = 1

        for activity, goal in itertools.zip_longest(new_content, goals, fillvalue=None):
            for slide in presentation.slides:
                for shape in slide.shapes:

                    activities = '{activity' + f'{num}' + "}"
                    description = '{description' + f'{num}' + "}"
                    deliverable = '{deliverable' + f'{num}' + "}"
                    spec_goal = '{specific_goal' + f'{num}' + "}"

                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.text == '{context}':
                                paragraph.text = method_data.context
                                paragraph.font.size = Pt(14)
                                paragraph.font.color.rgb = RGBColor(84, 83, 84)
                            elif paragraph.text == '{main_goal}':
                                paragraph.text = method_data.main_goal
                                paragraph.font.size = Pt(14)
                                paragraph.font.color.rgb = RGBColor(255, 255, 255)

                            for run in paragraph.runs:
                                if activity is not None:
                                    if run.text == activities:
                                        run.text = activity.activity.name
                                        run.font.size = Pt(12)
                                        run.font.color.rgb = RGBColor(0, 0, 0)
                                    if run.text == description:
                                        run.text = activity.description
                                        run.font.size = Pt(12)
                                        run.font.color.rgb = RGBColor(0, 0, 0)
                                    if run.text == deliverable:
                                        run.text = activity.deliverables
                                        run.font.size = Pt(12)
                                        run.font.color.rgb = RGBColor(84, 83, 84)
                                if goal is not None:
                                    if run.text == spec_goal:
                                        run.text = goal
                                        run.font.size = Pt(14)
                                        run.font.color.rgb = RGBColor(255, 255, 255)

            num += 1

        print("Methodology File modified successfully.")

    except Exception as e:
        raise Exception(f"Replace information error: {e}")


def include_proposal_info_ppt(presentation, proposal_id):
    '''
    Function to include Budget and general info from Proposal (Use cases)
    into Sharepoint template in Memory.
    '''

    try:
        # Return the saved Budget information from Proposal Instance Model
        proposal_info = ProposalsModel.objects.get(id=proposal_id)
        data_proposal = proposal_info.budget_info

        # Extract investment values
        key, value = data_proposal.popitem()

        # Return use cases
        use_cases = UseCasesModel.objects.filter(proposal=proposal_id)

        data = {
                '{project_price}': value['price'], '{duration}': value['duration'],
                '{considerations}': value['considerations'],
                '{solution_name}': proposal_info.solution.name if proposal_info.solution else '',
            }

        if use_cases:

            for num, use_case in enumerate(use_cases, start=1):

                temp_dict = {
                    '{usecase_name' + f"{num}" + '}': use_case.name,
                    '{usecase_description' + f"{num}" + '}': use_case.description,
                    '{usecase_situation' + f"{num}" + '}': use_case.situation,
                    '{usecase_task' + f"{num}" + '}': use_case.task,
                    '{usecase_actions' + f"{num}" + '}': use_case.action,
                    '{usecase_results' + f"{num}" + '}': use_case.results
                }

                data.update(temp_dict)

        search_and_include_info(presentation, data)

        print("Proposal file modified successfully.")

    except Exception as e:
        raise Exception(f"Replace information error: {e}")