# Office Libraries
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Django
from django.conf import settings


# Dataframes
import pandas as pd

# import matplotlib
# matplotlib.use('Agg')  # Set the backend to Agg before importing pyplot

# import matplotlib.pyplot as plt

# Memory
from io import BytesIO


def include_excel_as_table(ppt_file, excel_file, sheet_name=0):
    '''
    Function to include the RFI Excel template into the PPTX Presentation

    ppt_file: File pptx in memory buffer to include RFI Excel
    excel_file: RFI Excel file in memory buffer
    '''

    # Read the Excel file into a DataFrame
    df = pd.read_excel(excel_file, sheet_name=sheet_name)

    # Load the PowerPoint presentation
    try:
        presentation = Presentation(ppt_file)
    except Exception as e:
        raise Exception(f"Failed to open presentation: {e}")

    # Add a new slide with a blank layout
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])

    # Define table size and position
    rows, cols = df.shape
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(11)
    height = Inches(5)

    try:
        # TODO: Customize the fonts of header if needed
        # slide.shapes.title.text = 'Request for Information'
        # slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(32)
        # slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(94, 63, 209)

        # Add table to the slide
        table = slide.shapes.add_table(rows+1, cols, left, top, width, height).table
    except Exception as e:
        raise Exception(f"The table couldn't be added to the slide: {e}")

    try:
        # Set column names in the header row
        for col_idx, col_name in enumerate(df.columns):
            cell = table.cell(0, col_idx)
            cell.text = str(col_name)
            # Set font size and color for the header
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(8)
                    run.font.color.rgb = RGBColor(255, 255, 255)
                paragraph.alignment = PP_ALIGN.CENTER

            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(68, 46, 142)

        # Add the DataFrame content to the table
        for row_idx, row in df.iterrows():
            for col_idx, value in enumerate(row):
                cell = table.cell(row_idx+1, col_idx)
                cell.text = str(value)
                # Set font size and color for the data cells
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(8)
                        run.font.color.rgb = RGBColor(0, 0, 0)
                    paragraph.alignment = PP_ALIGN.CENTER

                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(215, 211, 227)

        print("The process to include the table completed succesfully")

    except Exception as e:
        raise Exception(f"Error mapping the dataframe: {e}")


    output_sp_file = BytesIO()
    presentation.save(output_sp_file)
    output_sp_file.seek(0)

    return output_sp_file

