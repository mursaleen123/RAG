# Standard Libraries
import requests

# Django
from django.conf import settings


# S3 Libraries
import boto3
from botocore.exceptions import NoCredentialsError, PartialCredentialsError

# Models
#from MVPautomation.sharepoint.utils.get_file_memory import get_sp_file
#from MVPautomation.users.models import User

# Office Libraries
from pptx.util import Inches, Pt

# Memory
from io import BytesIO

# Utils
from .delete_functions import get_slide_index


def include_logo(presentation, logo_path):
    '''
    Function to include the company logo in the ppt file
    '''

    ppt_file = get_sp_file(logo_path)

    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.name == "Logo":
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height

                sp = shape._element
                sp.getparent().remove(sp)

                slide.shapes.add_picture(ppt_file, left, top, width, height)


def download_image_from_s3(s3_key_path):
    '''
    Downloads a file from an S3 bucket.

    bucket_name: Name of the S3 bucket.
    s3_key: Key (path) of the file in the S3 bucket.
    local_file_name: Path to save the file locally.
    '''

    bucket_name = 'rap-storage'

    # Initialize the S3 client
    s3_client = boto3.client('s3')

    try:
        # Download the file from S3 to local file system
        image_stream = BytesIO()
        image_s3 = s3_client.download_fileobj(bucket_name, s3_key_path, image_stream)
        image_stream.seek(0)

        print(f"File {s3_key_path} downloaded successfully.")
    except Exception as e:
        print(f"An error occurred: {e}")

    return image_stream


def add_images_s3_ppt(presentation, data):
    '''
    Function to include images from S3 into a ppt file
    '''

    data = data['introduce_team']

    flag = 1

    try:
        for user in data:
            for slide in presentation.slides:
                for shape in slide.shapes:
                    members = ['{member' + f'{flag}' + "_photo}"]

                    if shape.has_text_frame:
                        for index, paragraph in enumerate(shape.text_frame.paragraphs):
                            if any(member in paragraph.text for member in members):
                                # Delete any space in text
                                text = paragraph.text.replace("\u200b", "")

                                # Get image path from S3
                                key_path = User.objects.get(id=user['{id}']).image
                                # Get image stream
                                image_stream = download_image_from_s3(str(key_path))


                                # Get slide index based on a keyword to identify the target slide
                                keyword = ["[Team-our-team]"]
                                slide_index = get_slide_index(presentation, keyword)
                                slide = presentation.slides[slide_index[0]]

                                # Specify the position and size of the image in Inches
                                left_pos = shape.left
                                top_pos = shape.top + index * Inches(0.3)

                                width = Inches(0.89)
                                height = Inches(0.89)

                                # Add the image to the slide    
                                try:
                                    slide.shapes.add_picture(image_stream, left_pos, top_pos, width, height)
                                    print("Image added succesfully")
                                except Exception as e:
                                    print(f"Error adding image: {e}")
                                    break

            flag += 1

    except Exception as e:
        raise Exception(f"Upload images error: {e}")