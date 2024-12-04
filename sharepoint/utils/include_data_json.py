# Standard Libraries

# Django
from django.conf import settings


# Django
from django.core.exceptions import ObjectDoesNotExist

# Models
#from MVPautomation.project.models import PresentationDataModel, ProposalVersionModel

# Office Libraries
from pptx.util import Pt


def get_all_values(data, values=set()):
    '''
    Get all the unique values in the JSON-like dict or list
    '''

    if isinstance(data, dict):
        for value in data.values():
            if isinstance(value, (dict, list)):
                get_all_values(value, values)
            else:
                values.add(value)

    elif isinstance(data, list):
        for item in data:
            get_all_values(item, values)

    return values


def get_all_keys(data, keys=set()):
    '''
    Get all the keys in the JSON dict
    '''

    if isinstance(data, dict):
        for key, value in data.items():
            keys.add(key)
            get_all_keys(value, keys)

    elif isinstance(data, list):
        for item in data:
            get_all_keys(item, keys)

    return keys


def find_key_in_json(data, key_to_find):
    '''
    Function to get the Value for a specific key in a JSON object
    '''
    if isinstance(data, dict):
        for key, value in data.items():
            if key == key_to_find:
                return {key: value}

            result = find_key_in_json(value, key_to_find)
            if result is not None:
                return result

    elif isinstance(data, list):
        for item in data:
            result = find_key_in_json(item, key_to_find)
            if result is not None:
                return result
    return None


def populate_versioning(presentation, instance_id):
    '''
    Function to include the JSON data stored in DB to replace keywords in slides presentation.

    - presentation: PPTX memory file
    - instance_id: ID of the ProposalVersionModel instance
    '''
    # Creates a empty list to include all the keys and values from the new version of the doc
    all_content = []

    data = ProposalVersionModel.objects.get(id=instance_id)

    all_content.extend([data.team, data.budget_info, data.methodology, data.use_cases])

    try:
        search_and_include_info(presentation, data)

        print("The process to include JSON in PPTx made successfully")
    except Exception as e:
        raise Exception(f"Mapping information in PPTX failed: {e}")


def include_json_into_ppt(presentation, content):
    '''
    Function to include the JSON data stored in DB to replace keywords in slides presentation.
    '''
    if isinstance(content, int):

        try:
            data = PresentationDataModel.objects.get(id=content).json_data
        except ObjectDoesNotExist:
            raise Exception(f"PresentationDataModel with id {content} does not exist")
    else:
        data = content

    try:
        search_and_include_info(presentation, data)
        print("The process to include JSON in PPTx completed successfully")
    except Exception as e:
        raise Exception(f"Mapping information in PPTX failed: {e}")


def search_and_include_info(presentation, data):
    '''
    Function to search the keywords in PTTX file and replace with the new information associated with it.

    - presentation: PPTX memory file
    - data: List of dict that will contain the keys and values to replace

    - Return: No data return
    '''

    # Group in a list all the keys available in the JSON
    keywords = get_all_keys(data)
    keywords = list(keywords)

    # Ensure the key_name is stripped of any zero-width space characters
    keywords = [kw.replace("\u200b", "").strip() for kw in keywords]

    for slide in presentation.slides:
        for shape in slide.shapes:

            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    # Check if any keyword in the list is in the run's text
                    if any(keyword in paragraph.text for keyword in keywords):
                        # Delete any space in text
                        text = paragraph.text.replace("\u200b", "")
                        # Get the value of the key found in the slide
                        values = find_key_in_json(data, text)

                        if values is not None:
                            key, value = values.popitem()
                            # Reeplace value in slide key
                            paragraph.text = str(value)
                            paragraph.font.size = Pt(12)
                            break
                        else:
                            print(f"Missing value for key: {paragraph.text}")
                            pass