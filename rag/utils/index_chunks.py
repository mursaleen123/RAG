# Standard Libraries
import requests

# Django Settings
from django.conf import settings

# Azure Search Libraries
from azure.search.documents import SearchClient

# Utils
# from MVPautomation.utils.utils_token import refresh_token, is_token_expired, get_actual_token

# Base-64 Key
import base64

# Power Point 
from pptx import Presentation



def get_file_content_from_sharepoint(file_url):
    '''
    Get the entire content of the file from Sharepoint

    - file_url: Entire string path to the file stored in Sharepoint

    Return:
    - The whole content of the file
    '''

    # if is_token_expired(organization='sharepoint') is True:
    #     refresh_token(organization='sharepoint')

    # token = get_actual_token(organization='sharepoint')

    # headers = {
    #     'Authorization': f'Bearer {token}',
    #     'Accept': 'application/json'
    # }

    # Request the chunk from SharePoint
    # relative_path = file_url[41:]
    # url = f"https://iotaimpactcompany1.sharepoint.com/sites/iotaimpact/_api/web/GetFileByServerRelativeUrl('{relative_path}')/$value"

    # response = requests.get(url, headers=headers)

    # if response.status_code == 200:
    #     print(f"Size: {len(response.content)} bytes")

    #     yield response.content

    # else:
    #     # Any other status code or issue
    #     raise Exception(f"Error fetching file: {response.status_code}, {response.text}")
    
    if file_url.lower().endswith(".pptx"):
        try:
            # Assuming 'file_url' points to a local file for simplicity
            with open(file_url, 'rb') as file:
                presentation = Presentation(file)
                content = []
                
                # Extract text from each slide
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text'):
                            content.append(shape.text)
                
                # Join all slide texts into one string
                return "\n".join(content)
        except FileNotFoundError:
            raise Exception(f"File not found at path: {file_url}")
        except Exception as e:
            raise Exception(f"Error extracting content from file: {str(e)}")
    
    else:
        # If not a PowerPoint file, proceed with the existing binary reading process
        try:
            with open(file_url, 'rb') as file:
                content = file.read()  # Read the entire content of the file
                return content  # Yield the content as a generator
        except FileNotFoundError:
            raise Exception(f"File not found at path: {file_url}")
        except Exception as e:
            raise Exception(f"Error reading file: {str(e)}")
    


def split_file_into_chunks(file_content, chunk_size=5 * 1024 * 1024):
    '''
    Splits the file content into chunks of a specified size.

    - file_content: The binary content of the file.
    - chunk_size: The size of each chunk in bytes (10 MB).

    Return:
    - Each chunk of binary file content.
    '''

    # Convert to list from a generator obj

    for i in range(0, len(file_content), chunk_size):
        yield file_content[i:i + chunk_size]


def index_chunks_to_search(file_metadata, credential, index_name):
    '''
    Function to classify the splitted chunks and uploaded in search client

    - file_metadata: Full content of the file selected
    - credential: Azure credentials
    - index_name: index name from RAG
    '''

    search_client = SearchClient(
        endpoint=settings.RAG_ENDPOINT, index_name=index_name, credential=credential
    )

    # Metadata to classify the chunks
    file_name = file_metadata['id']
    file_url = file_metadata['webUrl']

    content = get_file_content_from_sharepoint(file_url)
    
    if content is None:
            raise ValueError(f"Failed to retrieve content for file: {file_url}") 
    
    if isinstance(content, bytes):
        content = content.decode('utf-8', errors='ignore')

    chunk_size = 5 * 1024 * 1024  # 5 MB chunks

    for chunk_number, chunk in enumerate(split_file_into_chunks(content, chunk_size), start=1):

        document_key = f"{file_name}_chunk_{chunk_number}"
        document_key_safe = base64.urlsafe_b64encode(document_key.encode()).decode()
        
        document = {
            "id": document_key_safe,  # Use the Base64-encoded key
            "metadata_spo_item_name": file_name,
            "content": chunk,
            "metadata_spo_item_last_modified": file_metadata['lastModifiedDateTime'],
            "metadata_spo_item_path": file_url,
            "metadata_spo_item_content_type": file_url,
            "metadata_spo_item_size": len(chunk),
        }

        search_client.upload_documents(documents=[document])

